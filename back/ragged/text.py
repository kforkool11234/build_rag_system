import os
import re
import torch
from typing import List
import PyPDF2
from pptx import Presentation
from docx import Document
import nltk
from nltk.tokenize import sent_tokenize
from transformers import pipeline
import open_clip
import chromadb
from tqdm.auto import tqdm
from decouple import config
# Download NLTK resources
nltk.download('punkt')
tokenizer = open_clip.get_tokenizer('ViT-B-32')
# Check for GPU and set device
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"Using device: {device}")

# Initialize models with GPU support if available
summarizer = pipeline(
    "summarization",
    model="sshleifer/distilbart-cnn-12-6",
    device=0 if device == "cuda" else -1  # Use GPU if available
)
embedding_model,_,_ = open_clip.create_model_and_transforms('ViT-B-32', pretrained='laion2b_s34b_b79k')

# Initialize Chroma DB
CHROMA_DB_URL=config('chroma_url')
chroma_client = chromadb.HttpClient(host=CHROMA_DB_URL)



# ----------------------------------------------------------------------------------
# File Processing Functions
# ----------------------------------------------------------------------------------
def extract_text_from_file(file_path: str) -> str:
    """Extract text from PDF, DOCX, PPTX, or TXT files."""
    text = ""
    
    if file_path.endswith('.pdf'):
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file_path.endswith('.pptx'):
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()

    else:
        raise ValueError("Unsupported file format")
    
    return text
def split_paragraphs(text: str) -> List[str]:
    """Split text into paragraphs."""
    paragraphs = re.split(r'\n\s*\n', text.strip())
    return [para.replace('\n', ' ').strip() for para in paragraphs if para.strip()]

def chunk_content(text: str, word_limit: int = 1000) -> List[str]:
    """Split text into chunks with sentence boundaries."""
    chunks = []
    current_chunk = []
    current_word_count = 0
    sentences = sent_tokenize(text)
    
    for sentence in sentences:
        words = sentence.split()
        sentence_word_count = len(words)
        
        if current_word_count + sentence_word_count > word_limit:
            if current_chunk:
                chunks.append(' '.join(current_chunk).strip())
                current_chunk = []
                current_word_count = 0
                
            if sentence_word_count > word_limit:
                partial = ' '.join(words[:word_limit]) + '...'
                chunks.append(partial)
                continue
                
        current_chunk.append(sentence)
        current_word_count += sentence_word_count
    
    if current_chunk:
        chunks.append(' '.join(current_chunk).strip())
    
    return chunks

def process_file(file_path: str) -> tuple:
    """Return chunks and filename."""
    full_text = extract_text_from_file(file_path)
    paragraphs = split_paragraphs(full_text)
    chunks = []
    for para in paragraphs:
        chunks.extend(chunk_content(para))
    return chunks, os.path.basename(file_path)

# ----------------------------------------------------------------------------------
# Summarization and Storage Functions
# ----------------------------------------------------------------------------------
def summarize_chunks_batch(chunks: List[str]) -> List[str]:
    """Summarize batches of chunks together for faster processing."""
    batch_size = 4  # Number of chunks to summarize together
    summaries = []
    
    with tqdm(total=len(chunks), desc="Summarizing chunks") as pbar:
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i+batch_size]
            combined_text = "\n".join(batch)
            
            summary = summarizer(
                combined_text,
                max_length=150,  # Increased max length for combined chunks
                min_length=50,
                do_sample=False,
                truncation=True
            )
            summaries.append(summary[0]['summary_text'].strip())
            pbar.update(len(batch))
    
    return summaries

def store_in_chroma(c_id, chunks: List[str], filename: str, mode: str = "summary"):
    collection = chroma_client.get_or_create_collection(
        name=c_id,
        metadata={"hnsw:space": "cosine"}
    )

    batch_size = 8 if device == "cuda" else 4

    final_ids = []
    final_metadatas = []
    final_documents = []
    final_embeddings = []

    if mode == "summary":
        group_size = 4  # Number of chunks per summary
        summaries = summarize_chunks_batch(chunks, group_size=group_size)  # Make sure this respects group size

        # Get one embedding per summary
        with torch.no_grad():
            tokens = tokenizer(summaries).to(device)
            group_embeddings = embedding_model.encode_text(tokens)
            group_embeddings = group_embeddings / group_embeddings.norm(dim=-1, keepdim=True)
            group_embeddings = group_embeddings.cpu().tolist()
        # Distribute group embeddings to each chunk in its group
        for i, group_embedding in enumerate(group_embeddings):
            start_idx = i * group_size
            end_idx = min(start_idx + group_size, len(chunks))
            for j in range(start_idx, end_idx):
                final_documents.append(chunks[j])
                final_embeddings.append(group_embedding)
                final_metadatas.append({
                    "filename": filename,
                    "mode": mode,
                    "summary": summaries[i],
                    "type": "text",
                    "original_chunk": chunks[j]
                })
                final_ids.append(f"{filename}_chunk_{j}")
    else:
        # No summarization: one embedding per chunk
        with torch.no_grad():
            tokens = tokenizer(chunks).to(device)
            final_embeddings = embedding_model.encode_text(tokens)
            final_embeddings = final_embeddings / final_embeddings.norm(dim=-1, keepdim=True)
            final_embeddings = final_embeddings.cpu().tolist()
        final_documents = chunks
        final_metadatas = [{
            "filename": filename,
            "mode": mode,
            "type": "text",
            "original_chunk": chunk
        } for chunk in chunks]
        final_ids = [f"{filename}_chunk_{i}" for i in range(len(chunks))]

    # Batch add to Chroma
    with tqdm(total=len(final_documents), desc="Storing in database") as pbar:
        for i in range(0, len(final_documents), 100):  # Chroma batch size
            collection.add(
                documents=final_documents[i:i + 100],
                embeddings=final_embeddings[i:i + 100],
                metadatas=final_metadatas[i:i + 100],
                ids=final_ids[i:i + 100]
            )
            pbar.update(min(100, len(final_documents) - i))

    print(f"\n✅ Stored {len(final_documents)} chunks from {filename} in mode: {mode}")



# ----------------------------------------------------------------------------------
# Main Workflow
# ----------------------------------------------------------------------------------
def text(file_name,c_id):    
    # Process file
    chunks, filename = process_file(file_name)
    print(f"\n📄 Processed {len(chunks)} chunks from {filename}")
    store_in_chroma(c_id,chunks, filename, mode="full")
    