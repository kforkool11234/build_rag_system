import os
import re
import torch
from typing import List
import PyPDF2
from pptx import Presentation
from docx import Document
import nltk
from nltk.tokenize import sent_tokenize
from transformers import pipeline
from sentence_transformers import SentenceTransformer
import chromadb
from tqdm.auto import tqdm

# Download NLTK resources
nltk.download('punkt')

# Check for GPU and set device
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"Using device: {device}")

# Initialize models with GPU support if available
summarizer = pipeline(
    "summarization",
    model="sshleifer/distilbart-cnn-12-6",
    device=0 if device == "cuda" else -1  # Use GPU if available
)
embedding_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2', device=device)

# Initialize Chroma DB
chroma_client = chromadb.PersistentClient(path="chroma_db")



# ----------------------------------------------------------------------------------
# File Processing Functions
# ----------------------------------------------------------------------------------
def extract_text_from_file(file_path: str) -> str:
    """Extract text from PDF, DOCX, PPTX, or TXT files."""
    text = ""
    
    if file_path.endswith('.pdf'):
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"

    elif file_path.endswith('.docx'):
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file_path.endswith('.pptx'):
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()

    else:
        raise ValueError("Unsupported file format")
    
    return text
def split_paragraphs(text: str) -> List[str]:
    """Split text into paragraphs."""
    paragraphs = re.split(r'\n\s*\n', text.strip())
    return [para.replace('\n', ' ').strip() for para in paragraphs if para.strip()]

def chunk_content(text: str, word_limit: int = 1000) -> List[str]:
    """Split text into chunks with sentence boundaries."""
    chunks = []
    current_chunk = []
    current_word_count = 0
    sentences = sent_tokenize(text)
    
    for sentence in sentences:
        words = sentence.split()
        sentence_word_count = len(words)
        
        if current_word_count + sentence_word_count > word_limit:
            if current_chunk:
                chunks.append(' '.join(current_chunk).strip())
                current_chunk = []
                current_word_count = 0
                
            if sentence_word_count > word_limit:
                partial = ' '.join(words[:word_limit]) + '...'
                chunks.append(partial)
                continue
                
        current_chunk.append(sentence)
        current_word_count += sentence_word_count
    
    if current_chunk:
        chunks.append(' '.join(current_chunk).strip())
    
    return chunks

def process_file(file_path: str) -> tuple:
    """Return chunks and filename."""
    full_text = extract_text_from_file(file_path)
    paragraphs = split_paragraphs(full_text)
    chunks = []
    for para in paragraphs:
        chunks.extend(chunk_content(para))
    return chunks, os.path.basename(file_path)

# ----------------------------------------------------------------------------------
# Summarization and Storage Functions
# ----------------------------------------------------------------------------------
def summarize_chunks_batch(chunks: List[str]) -> List[str]:
    """Summarize batches of chunks together for faster processing."""
    batch_size = 4  # Number of chunks to summarize together
    summaries = []
    
    with tqdm(total=len(chunks), desc="Summarizing chunks") as pbar:
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i+batch_size]
            combined_text = "\n".join(batch)
            
            summary = summarizer(
                combined_text,
                max_length=150,  # Increased max length for combined chunks
                min_length=50,
                do_sample=False,
                truncation=True
            )
            summaries.append(summary[0]['summary_text'].strip())
            pbar.update(len(batch))
    
    return summaries

def store_in_chroma(c_id,chunks: List[str], filename: str, mode: str = "summary"):
    collection = chroma_client.get_or_create_collection(
    name=c_id,  # Unified name
    metadata={"hnsw:space": "cosine"}
)
    """Store chunks with optional summarization."""
    batch_size = 8 if device == "cuda" else 4

    if mode == "summary":
        # Summarize chunks in batches (grouping 4 chunks at a time)
        summaries = summarize_chunks_batch(chunks)
        # Compute one embedding per summary (each summarizing 4 chunks)
        group_embeddings = embedding_model.encode(summaries, batch_size=batch_size).tolist()

        # Now, for each group of 4 chunks, duplicate the group embedding for each chunk
        final_ids = []
        final_metadatas = []
        final_documents = []
        final_embeddings = []
        group_size = 4  # Number of chunks per summary
        num_groups = len(summaries)  # Should be ceil(len(chunks) / group_size)
        for i in range(num_groups):
            start_idx = i * group_size
            end_idx = min(start_idx + group_size, len(chunks))
            # For each chunk in the current group, assign the same embedding (from group_embeddings[i])
            for j in range(start_idx, end_idx):
                final_ids.append(f"{filename}_chunk_{j}")
                final_metadatas.append({"filename": filename, "mode": mode,"type": "text", "original_chunk": chunks[j]})
                final_documents.append(chunks[j])
                final_embeddings.append(group_embeddings[i])

        # For debugging, you can uncomment the next line to ensure all lists have the same length.
        # assert len(final_ids) == len(final_metadatas) == len(final_documents) == len(final_embeddings)
    else:
        # Use full chunks without summarization
        final_embeddings = embedding_model.encode(chunks, batch_size=batch_size).tolist()
        final_documents = chunks
        final_metadatas = [{"filename": filename, "mode": mode,"type":"text", "original_chunk": chunk} for chunk in chunks]
        final_ids = [f"{filename}_chunk_{i}" for i in range(len(chunks))]

    # Batch add to Chroma
    with tqdm(total=len(final_documents), desc="Storing in database") as pbar:
        for i in range(0, len(final_documents), 100):  # Chroma batch size
            collection.add(
                documents=final_documents[i:i+100],
                embeddings=final_embeddings[i:i+100],
                metadatas=final_metadatas[i:i+100],
                ids=final_ids[i:i+100]
            )
            pbar.update(min(100, len(final_documents) - i))

    print(f"\n✅ Stored {len(final_documents)} chunks from {filename} in mode: {mode}")



# ----------------------------------------------------------------------------------
# Main Workflow
# ----------------------------------------------------------------------------------
def text(file_name,c_id):    
    # Process file
    chunks, filename = process_file(file_name)
    print(f"\n📄 Processed {len(chunks)} chunks from {filename}")
    store_in_chroma(c_id,chunks, filename, mode="full")
    